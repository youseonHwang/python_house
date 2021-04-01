from flask import Flask, render_template, request, redirect, url_for
from pptx import Presentation  # 필요한 라이브러리 import
from pptx.util import Inches  # 사진, 표 등을 그리기 위해

import os
import database

import sys
application = Flask(__name__)


@application.route("/")
def hello():
    return render_template("hello.html")


@application.route("/apply")
def apply():
    return render_template("apply.html")


@application.route("/applyphoto")
def photo_apply():
    location = request.args.get("location")
    cleaness = request.args.get("clean")
    built_in = request.args.get("built")
    if cleaness is None:
        cleaness = False
    else:
        cleaness = True
    database.save(location, cleaness, built_in)
    return render_template("apply_photo.html")


@application.route("/upload_done", methods=["POST"])
def upload_done():
    uploaded_files = request.files["file"]
    uploaded_files.save("static/img/{}.jpeg".format(database.now_index()))
    return redirect(url_for("hello"))


@application.route("/list")
def list():
    house_list = database.load_list()
    length = len(house_list)  # html에서 length를 사용하기 위해서 받아서 넘김
    return render_template("list.html", house_list=house_list, length=length)

# house의 세부내용 보는 함수


@application.route("/house_info/<int:index>/")
def house_info(index):
    house_info = database.load_house(index)
    print(house_info)
    location = house_info["location"]
    cleaness = house_info["cleaness"]
    built_in = house_info["built_in"]
    photo = f"img/{index}.jpeg"  # static안의 img안에 index로 이미지 가져오기
    return render_template("house_info.html",
                           location=location, cleaness=cleaness,
                           built_in=built_in, photo=photo)

# -------------------------------------------------------------------------------------------------------------------------------


@application.route("/addSlide")
def add_slide():
    print("addSlide도착")
    prs = Presentation("static/yescnc_ppt_example.pptx")
    prs.slides.add_slide(prs.slide_layouts[2])  # 빈 슬라이드 추가
    prs.save("static/yescnc_ppt_example.pptx")
    return redirect(url_for("main"))


@application.route("/main")
def main():
    prs = Presentation("static/yescnc_ppt_example.pptx")
    print(prs)

    return render_template("main.html", slides=prs.slides)


@application.route("/editSlide", methods=["POST"])
def edit_slide():
    # 프레젠테이션 객체 생성
    prs = Presentation("static/yescnc_ppt_example.pptx")

    # 슬라이드 마스터 객체 생성
    master = prs.slide_master

    # reqeust에 실린 내용 확인
    for key, value in request.form.items():
        print(key, value)

    # 슬라이드 수만큼 반복문 돌기
    for slide_index in range(0, len(prs.slides)):

        # 0번째 슬라이드부터 시작
        slide = prs.slides[slide_index]
        print("--------[%d] ------ " % (slide_index))

        # 현재 슬라이드의 레이아웃 확인하기
        if master.slide_layouts[0] == slide.slide_layout:  # 제목 슬라이드라면
            print(slide.slide_layout.name)
            print("0번째 슬라이드 맞음")

            # 슬라이드 안에 있는 모양들 반복문 돌기
            for shape_index, shape in enumerate(slide.shapes):

                # 현재 모양 확인
                print("shape의 이름: "+(shape.name))

                # 모든 모양에 텍스트 프레임이 있는 것은 아니니깐 texxt프레임 있는지 확인
                if not shape.has_text_frame:
                    continue
                print("shape_index는 "+str(shape_index))
                if shape.name == "제목 1":  # 제목이라면

                    print("shape의 text: "+(shape.text_frame.text))
                    title = request.form.get("title"+str(slide_index+1))
                    print(title)

                    # 현재 텍스트 내용 지우기
                    shape.text_frame.clear()
                    print("!!textframe 클리어!!")

                    # 텍스트 내용 바꾸기
                    shape.text_frame.text = title
                    print("바꾼 결과:"+(shape.text_frame.text))

                if shape.name == "Text Placeholder 12" and shape_index == 3:
                    print("shape의 text: "+(shape.text_frame.text))
                    writer = request.form.get("writer"+str(slide_index+1))
                    print(writer)

                    # 현재 텍스트 내용 지우기
                    shape.text_frame.clear()
                    print("!!textframe 클리어!!")

                    # 텍스트 내용 바꾸기
                    shape.text_frame.text = writer
                    print("바꾼 결과:"+(shape.text_frame.text))

                if shape.name == "Text Placeholder 12" and shape_index == 4:
                    print("shape의 text: "+(shape.text_frame.text))
                    date = request.form.get("date"+str(slide_index+1))
                    print(date)

                    # 현재 텍스트 내용 지우기
                    shape.text_frame.clear()
                    print("!!textframe 클리어!!")

                    # 텍스트 내용 바꾸기
                    shape.text_frame.text = date
                    print("바꾼 결과:"+(shape.text_frame.text))

        if master.slide_layouts[1] == slide.slide_layout:  # 제목 및 내용 슬라이드라면
            print(slide.slide_layout.name)
            print("1번째 슬라이드 맞음")

            # 슬라이드 안에 있는 모양들 반복문 돌기
            for shape_index, shape in enumerate(slide.shapes):

                # 현재 모양 확인
                print("shape의 이름: "+(shape.name))

                # 모든 모양에 텍스트 프레임이 있는 것은 아니니깐 texxt프레임 있는지 확인
                if not shape.has_text_frame:
                    continue

                print("shape_index는 "+str(shape_index))

                if shape.name == "Title 1":  # 제목이라면

                    print("shape의 text: "+(shape.text_frame.text))
                    title = request.form.get("title"+str(slide_index+1))
                    print(title)

                    # 현재 텍스트 내용 지우기
                    shape.text_frame.clear()
                    print("!!textframe 클리어!!")

                    # 텍스트 내용 바꾸기
                    shape.text_frame.text = title
                    print("바꾼 결과:"+(shape.text_frame.text))

                if shape.name == "TextBox 3":
                    print("shape의 text: "+(shape.text_frame.text))
                    content = request.form.get("content"+str(slide_index+1))
                    print(content)

                    # 현재 텍스트 내용 지우기
                    shape.text_frame.clear()
                    print("!!textframe 클리어!!")

                    # 텍스트 내용 바꾸기
                    shape.text_frame.text = content
                    print("바꾼 결과:"+(shape.text_frame.text))

        if master.slide_layouts[2] == slide.slide_layout:  # 사진 슬라이드라면
            print(slide.slide_layout.name)
            print("2번째 슬라이드 맞음")

            # 슬라이드 안에 있는 모양들 반복문 돌기
            for shape_index, shape in enumerate(slide.shapes):

                # 현재 모양 확인
                print("shape의 이름: "+(shape.name))

                # 모든 모양에 텍스트 프레임이 있는 것은 아니니깐 texxt프레임 있는지 확인
                if not shape.has_text_frame:
                    continue

                print("shape_index는 "+str(shape_index))

                if shape.name == "제목 3":  # 제목이라면

                    print("shape의 text: "+(shape.text_frame.text))
                    title = request.form.get("title"+str(slide_index+1))
                    print(title)

                    # 현재 텍스트 내용 지우기
                    shape.text_frame.clear()
                    print("!!textframe 클리어!!")

                    # 텍스트 내용 바꾸기
                    shape.text_frame.text = title
                    print("바꾼 결과:"+(shape.text_frame.text))
                if shape.name == "내용 개체 틀 4":
                    print("shape의 text: "+(shape.text_frame.text))
                    photo_content = request.form.get(
                        "photo_content"+str(slide_index+1))
                    print(photo_content)

                    # 현재 텍스트 내용 지우기
                    shape.text_frame.clear()
                    print("!!textframe 클리어!!")

                    # 텍스트 내용 바꾸기
                    shape.text_frame.text = photo_content
                    print("바꾼 결과:"+(shape.text_frame.text))

            # 만약 파일이 존재한다면
            file = ("photo_file"+str(slide_index+1))
            if request.files[file] is not None:
                print("사진이 존재합니다.")
                file_name = request.files[file].filename
                print("file_name은" + file_name)
                # extends_index = file.find(".")
                # file_name = file_name[:extends_index]
                # print(file_name)
                # print(request.files[file])

                uploaded_files = request.files[file]

                # 가져온 파일 저장
                uploaded_files.save("static/img/{}.jpeg".format(file_name))
                img_path = "static/img/{}.jpeg".format(file_name)
                left = top = Inches(1)
                width = height = Inches(1)
                # width, hegith가 없을 경우 원본 사이즈로
                slide.shapes.add_picture(
                    img_path, left, top, width=width, height=height)
    prs.save("static/new_ppt.pptx")
    return redirect(url_for("hello"))


if __name__ == "__main__":
    application.run(host='0.0.0.0')
